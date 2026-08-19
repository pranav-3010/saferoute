import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_enhanced_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Color Palette
    DARK_BLUE = RGBColor(15, 23, 42)       # #0f172a
    SKY_BLUE = RGBColor(2, 132, 199)       # #0284c7
    LIGHT_BLUE = RGBColor(224, 242, 254)   # #e0f2fe
    EMERALD_GREEN = RGBColor(16, 185, 129) # #10b981
    LIGHT_GREEN = RGBColor(236, 253, 245)  # #ecfdf5
    CRIMSON_RED = RGBColor(239, 68, 68)    # #ef4444
    LIGHT_RED = RGBColor(254, 242, 242)    # #fef2f2
    LIGHT_BG = RGBColor(248, 250, 252)     # #f8fafc
    DARK_TEXT = RGBColor(30, 41, 59)       # #1e293b
    MUTED_TEXT = RGBColor(100, 116, 139)   # #64748b
    CARD_BG = RGBColor(248, 250, 252)      # #f8fafc
    WHITE = RGBColor(255, 255, 255)
    BORDER_GRAY = RGBColor(226, 232, 240)  # #e2e8f0

    def add_slide_header(slide, title_text, slide_num_text):
        # Header Box
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(11.733), Inches(0.7))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{slide_num_text}  |  {title_text.upper()}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE

        # Underline bar
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.05), Inches(11.733), Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SKY_BLUE
        shape.line.fill.background()

    def add_flowchart_box(slide, left, top, width, height, text, bg_color, border_color, text_color=DARK_BLUE, font_size=12, bold=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        return shape

    def add_arrow(slide, left, top, width, height, color=SKY_BLUE):
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # ================= SLIDE 1: TITLE PAGE =================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Background Accent Shape
    bg_accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bg_accent.fill.solid()
    bg_accent.fill.fore_color.rgb = SKY_BLUE
    bg_accent.line.fill.background()

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
    p = t_box.text_frame.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2024"
    p.font.name = "Segoe UI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.333), Inches(0.5))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "TITLE PAGE"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Main Card Box
    rect = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.1), Inches(10.333), Inches(4.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = SKY_BLUE
    rect.line.width = Pt(2)

    tf = rect.text_frame
    tf.word_wrap = True
    
    fields = [
        ("• Problem Statement ID — ", "SIH1642"),
        ("• Problem Statement Title — ", "AI-Powered Women Safety Navigation & Emergency Response System"),
        ("• Theme — ", "Women Safety / Smart Vehicles / AI-GIS Navigation"),
        ("• PS Category — ", "Software"),
        ("• Team ID — ", "[Enter Your Registered Team ID]"),
        ("• Team Name (Registered on portal) — ", "[Enter Your Registered Team Name]")
    ]

    for i, (label, val) in enumerate(fields):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run1 = p.add_run()
        run1.text = label
        run1.font.name = "Segoe UI"
        run1.font.size = Pt(16)
        run1.font.bold = True
        run1.font.color.rgb = DARK_BLUE

        run2 = p.add_run()
        run2.text = val
        run2.font.name = "Segoe UI"
        run2.font.size = Pt(16)
        run2.font.bold = True
        run2.font.color.rgb = SKY_BLUE

    # ================= SLIDE 2: PROBLEM STATEMENT (WITH DIAGRAM) =================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide2, "1. Problem Statement", "SLIDE 02")

    # Left Column Text
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.6))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Critical Flaw in Standard Navigation"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(10)

    points = [
        "• Shortest-Distance Trap: Standard mapping engines (Google Maps, Waze) optimize exclusively for distance & speed, ignoring safety.",
        "• Unlit Back-Alley Risk: Pedestrians are routed through dark, unlit side-streets to save 60 seconds of travel time.",
        "• Zero Safety Context: No real-time evaluation of street lighting, police station proximity, or crime hazard pins.",
        "• Mobile SOS Latency: Traditional panic apps suffer from mobile browser CORS restrictions during emergency situations."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)

    # Right Column: Visual Flowchart Diagram
    right_card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.6))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = BORDER_GRAY
    
    t_box = slide2.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.3), Inches(0.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = "Standard Maps vs. Safety Failure Diagram"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Flowchart Nodes (Vertical)
    add_flowchart_box(slide2, Inches(7.5), Inches(2.2), Inches(4.3), Inches(0.7), "Standard Navigation Engine\n(Optimizes Distance/Time Only)", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 12)
    add_arrow(slide2, Inches(9.4), Inches(3.0), Inches(0.5), Inches(0.4), CRIMSON_RED)
    
    add_flowchart_box(slide2, Inches(7.5), Inches(3.5), Inches(4.3), Inches(0.7), "Routes Pedestrian Through Dark Unlit Alley\n(To Save 60 Seconds)", LIGHT_RED, CRIMSON_RED, CRIMSON_RED, 12)
    add_arrow(slide2, Inches(9.4), Inches(4.3), Inches(0.5), Inches(0.4), CRIMSON_RED)

    add_flowchart_box(slide2, Inches(7.5), Inches(4.8), Inches(4.3), Inches(0.8), "❌ High Crime Exposure | ❌ Zero Streetlights\n❌ Zero Police Proximity Scoring", LIGHT_RED, CRIMSON_RED, CRIMSON_RED, 11)

    # ================= SLIDE 3: EXISTING SITUATION =================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide3, "2. Existing Situation & Competitor Gaps", "SLIDE 03")

    # Competitor Comparison Visual Cards (3 Columns)
    cards = [
        ("Google / Apple Maps", "🔴 Safety Rating: 0/5", "• Optimizes travel speed only\n• Ignores unlit dark streets\n• No police proximity scoring", LIGHT_RED, CRIMSON_RED),
        ("Standalone SOS Apps", "🟡 Safety Rating: 2/5", "• Reactive only (after danger)\n• No preventive safe routing\n• High user panic friction", LIGHT_BLUE, SKY_BLUE),
        ("Manual Panic Buttons", "🟡 Safety Rating: 2/5", "• Browser CORS latency\n• Requires manual phone press\n• Fails in fast emergencies", LIGHT_BLUE, SKY_BLUE)
    ]

    for i, (title, sub, body, bg_col, bdr_col) in enumerate(cards):
        left_pos = Inches(0.8 + i * 3.95)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.4), Inches(3.8), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_col
        shape.line.color.rgb = bdr_col
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = bdr_col
        p2.space_after = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(12)
        p3.font.color.rgb = DARK_TEXT

    # Market Opportunity Banner at Bottom
    bot_card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.733), Inches(1.5))
    bot_card.fill.solid()
    bot_card.fill.fore_color.rgb = LIGHT_GREEN
    bot_card.line.color.rgb = EMERALD_GREEN
    bot_card.line.width = Pt(1.5)

    tf = bot_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🟢 SafeRoute Market Opportunity: Preventive AI Safety Navigation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 95, 70)
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = "SafeRoute bridges the market gap by combining preventive AI safe routing (guiding women away from danger before it occurs) with zero-delay CORS-free emergency SOS automated cloud phone calls."
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_TEXT

    # ================= SLIDE 4: PROPOSED SOLUTION (5-FACTOR ENGINE DIAGRAM) =================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide4, "3. Proposed Solution — SafeRoute", "SLIDE 04")

    # Left Text
    left_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.4), Inches(5.6))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SafeRoute: 5-Factor AI Safety Engine"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(8)

    points = [
        "• 5-Factor Safety Scoring Model: Evaluates road geometry against 5 safety parameters (0 to 100).",
        "• 🟢 84/100 Safest Route Highlight: Renders recommended safe path in glowing emerald green.",
        "• 🔴 38/100 High-Risk Route Warning: Flags dark shortcuts and unlit side-streets in red.",
        "• Spherical Haversine Math: Sub-second spatial distance calculation for police stations and hospitals.",
        "• AI Explainability Cards: Tells users exactly WHY a route is recommended."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    # Right Column: Visual 5-Factor Diagram
    right_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.3), Inches(6.033), Inches(5.6))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = SKY_BLUE

    t_box = slide4.shapes.add_textbox(Inches(6.7), Inches(1.45), Inches(5.6), Inches(0.4))
    p = t_box.text_frame.paragraphs[0]
    p.text = "5-Factor Mathematical Safety Algorithm"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    factors = [
        ("💡 Street Lighting (30%)", "OSM tags & municipal lighting hierarchy"),
        ("🚓 Police Proximity (25%)", "Haversine distance <= 1.0 km radius"),
        ("⚠️ Hazard Incident Pins (20%)", "Snatching/harassment pins <= 500m"),
        ("🌙 Night Shift (15%)", "1.4x hazard penalty shift (8PM-5AM)"),
        ("🏪 Commercial Footfall (10%)", "Active retail density 'Eyes on Street'")
    ]

    for j, (title, desc) in enumerate(factors):
        add_flowchart_box(slide4, Inches(6.8), Inches(1.95 + j * 0.8), Inches(5.4), Inches(0.65), f"{title}\n{desc}", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 11)

    # Output Box at Bottom Right
    add_flowchart_box(slide4, Inches(6.8), Inches(6.05), Inches(5.4), Inches(0.65), "🟢 Output: 84/100 Safest Route Highlighted", LIGHT_GREEN, EMERALD_GREEN, RGBColor(6, 95, 70), 12)

    # ================= SLIDE 5: ARCHITECTURE & WORKFLOW DIAGRAM =================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide5, "4. Architecture & Data Workflow", "SLIDE 05")

    # Top Flowchart: System Architecture Pipeline
    t_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.4))
    p = t_box.text_frame.paragraphs[0]
    p.text = "System Navigation Architecture Workflow"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE

    # 4 Sequential Nodes
    add_flowchart_box(slide5, Inches(0.8), Inches(1.7), Inches(2.5), Inches(1.1), "1. User Query\n(Origin & Destination)", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 12)
    add_arrow(slide5, Inches(3.45), Inches(2.05), Inches(0.4), Inches(0.4))

    add_flowchart_box(slide5, Inches(4.0), Inches(1.7), Inches(2.5), Inches(1.1), "2. Nominatim Geocoder\n& OSRM Geometry", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 12)
    add_arrow(slide5, Inches(6.65), Inches(2.05), Inches(0.4), Inches(0.4))

    add_flowchart_box(slide5, Inches(7.2), Inches(1.7), Inches(2.5), Inches(1.1), "3. Haversine Math\n& 5-Factor Scoring", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 12)
    add_arrow(slide5, Inches(9.85), Inches(2.05), Inches(0.4), Inches(0.4))

    add_flowchart_box(slide5, Inches(10.4), Inches(1.7), Inches(2.133), Inches(1.1), "4. Leaflet Map\n🟢 Green / 🔴 Red", LIGHT_GREEN, EMERALD_GREEN, RGBColor(6, 95, 70), 12)

    # Bottom Flowchart: Emergency SOS Workflow Pipeline
    t_box2 = slide5.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.733), Inches(0.4))
    p = t_box2.text_frame.paragraphs[0]
    p.text = "Zero-Delay Emergency SOS Cloud Pipeline Workflow"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = CRIMSON_RED

    # 4 SOS Sequential Nodes
    add_flowchart_box(slide5, Inches(0.8), Inches(3.6), Inches(2.5), Inches(1.1), "User Taps SOS / Speaks\n'HELP ME' or 'BACHAO'", LIGHT_RED, CRIMSON_RED, CRIMSON_RED, 12)
    add_arrow(slide5, Inches(3.45), Inches(3.95), Inches(0.4), Inches(0.4), CRIMSON_RED)

    add_flowchart_box(slide5, Inches(4.0), Inches(3.6), Inches(2.5), Inches(1.1), "Native CORS-Free Anchor\n(Bypasses Mobile Latency)", LIGHT_RED, CRIMSON_RED, CRIMSON_RED, 12)
    add_arrow(slide5, Inches(6.65), Inches(3.95), Inches(0.4), Inches(0.4), CRIMSON_RED)

    add_flowchart_box(slide5, Inches(7.2), Inches(3.6), Inches(2.5), Inches(1.1), "n8n Cloud Webhook\n& Twilio REST API Node", LIGHT_RED, CRIMSON_RED, CRIMSON_RED, 12)
    add_arrow(slide5, Inches(9.85), Inches(3.95), Inches(0.4), Inches(0.4), CRIMSON_RED)

    add_flowchart_box(slide5, Inches(10.4), Inches(3.6), Inches(2.133), Inches(1.1), "Live Voice Call Placed\n+ Live GPS Broadcast", LIGHT_GREEN, EMERALD_GREEN, RGBColor(6, 95, 70), 12)

    # Technical Highlight Card
    tech_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.8))
    tech_card.fill.solid()
    tech_card.fill.fore_color.rgb = CARD_BG
    tech_card.line.color.rgb = BORDER_GRAY
    tf = tech_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Architectural Advantage: CORS-Free Zero-Delay Trigger"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = "Mobile browsers enforce strict CORS preflights on async fetch requests, delaying SOS dispatches by 3-6 seconds. SafeRoute uses native HTML anchor triggers targeting a hidden background iframe—guaranteeing 0ms execution without browser security blocks."
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = DARK_TEXT

    # ================= SLIDE 6: TECHNOLOGIES USED =================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide6, "5. Technologies Used", "SLIDE 06")

    # Table
    table_shape = slide6.shapes.add_table(9, 3, Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.6))
    table = table_shape.table
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(5.433)

    headers = ["Layer", "Technology Selected", "Technical Purpose / Value Provided"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(12)

    tech_data = [
        ("Frontend Core", "Vite + Vanilla ES6+ JavaScript", "Ultra-fast mobile rendering & zero-latency module loading"),
        ("Map Engine", "Leaflet JS Canvas Engine", "Interactive vector tile rendering & smooth polyline overlays"),
        ("Routing API", "OSRM (Open Source Routing)", "Real road network graph sampling & GeoJSON route geometry"),
        ("Geocoding", "Nominatim OSM Geocoder", "Address-to-GPS spatial coordinate resolution"),
        ("Spatial Engine", "Native Haversine Formula", "Spherical distance calculations for police & hazard scanning"),
        ("Automation", "n8n Workflow Automation", "Zero-delay webhook routing & automated cloud workflows"),
        ("Telephony Dispatch", "Twilio REST Voice API", "Instant automated phone call placement to emergency contacts"),
        ("Voice SOS Engine", "Web Speech Recognition API", "Hands-free panic phrase detection ('HELP ME', 'BACHAO')")
    ]

    for i, (l, t, p_text) in enumerate(tech_data):
        c0 = table.cell(i+1, 0)
        c1 = table.cell(i+1, 1)
        c2 = table.cell(i+1, 2)
        for c in [c0, c1, c2]:
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_BG
        c0.text_frame.paragraphs[0].text = l
        c1.text_frame.paragraphs[0].text = t
        c2.text_frame.paragraphs[0].text = p_text
        for c in [c0, c1, c2]:
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT

    # ================= SLIDE 7: PROTOTYPE DEMONSTRATION =================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide7, "6. Prototype Demonstration", "SLIDE 07")

    # Left Column
    left_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.6))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Live Prototype Capabilities"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(10)

    points = [
        "• Live Deployed Web App: Fully accessible live at saferoute-tawny.vercel.app",
        "• Multi-Modal Transit: Supports Car, Bike, Auto, and Walking routes.",
        "• 90+ Hyderabad Emergency Hubs: Pre-loaded with verified Police Stations, Hospitals, and Metro Hubs.",
        "• GPS Navigation Simulator: Animates live position movement along routes.",
        "• Off-Route 200m Safety Fence: Automatically detects deviations >200m off-route and pops up an emergency warning banner."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    # Right Column: Interface Diagram
    right_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.6))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = SKY_BLUE

    t_box = slide7.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.3), Inches(0.4))
    p = t_box.text_frame.paragraphs[0]
    p.text = "Key UI Components Demonstrated"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    add_flowchart_box(slide7, Inches(7.1), Inches(2.1), Inches(5.1), Inches(1.0), "Interface 1: Trip Input\n• Address Search & Geocoding\n• Mode Pills (Car/Bike/Auto/Walk)\n• Time Picker (Live Now / Night)", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 11)
    
    add_arrow(slide7, Inches(9.4), Inches(3.2), Inches(0.5), Inches(0.35))

    add_flowchart_box(slide7, Inches(7.1), Inches(3.65), Inches(5.1), Inches(1.4), "Interface 2: Interactive Route Map\n• 🟢 84/100 Safest Route Polyline\n• 🔴 38/100 High-Risk Route Warning\n• 🚓 90+ Police & Hospital Markers\n• 🚨 CORS-Free ONE-TAP SOS Button", LIGHT_GREEN, EMERALD_GREEN, RGBColor(6, 95, 70), 11)

    add_flowchart_box(slide7, Inches(7.1), Inches(5.2), Inches(5.1), Inches(1.2), "Navigation Features & Simulator\n• GPS Navigation Path Animation\n• 200m Off-Route Deviation Banner\n• Multilingual Voice SOS ('HELP ME')", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 11)

    # ================= SLIDE 8: INNOVATION & DIFFERENTIATORS =================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide8, "7. Innovation & Key Differentiators", "SLIDE 08")

    table_shape = slide8.shapes.add_table(7, 3, Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.733)
    table.columns[1].width = Inches(4.0)
    table.columns[2].width = Inches(4.0)

    headers = ["Feature / Capability", "Standard Navigation Apps", "SafeRoute AI Engine"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(12)

    diff_data = [
        ("Street Lighting Scoring", "❌ None (Time optimization only)", "🟢 3-Layer Hybrid Model"),
        ("Police Station Proximity", "❌ Ignored in route scoring", "🟢 Haversine Radius (d <= 1.0 km)"),
        ("Night Multiplier (8PM-5AM)", "❌ None", "🟢 1.4x Hazard Penalty Shift"),
        ("Commercial Footfall Bonus", "❌ None", "🟢 +10% 'Eyes on the Street'"),
        ("Mobile CORS-Free SOS", "❌ N/A", "🟢 0ms n8n + Twilio Voice Call"),
        ("Multilingual Voice SOS", "❌ None", "🟢 Hands-Free Web Speech API")
    ]

    for i, (f, s, sr) in enumerate(diff_data):
        c0 = table.cell(i+1, 0)
        c1 = table.cell(i+1, 1)
        c2 = table.cell(i+1, 2)
        for c in [c0, c1, c2]:
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_BG
        c0.text_frame.paragraphs[0].text = f
        c1.text_frame.paragraphs[0].text = s
        c2.text_frame.paragraphs[0].text = sr
        for c in [c0, c1, c2]:
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(11.5)
            p.font.color.rgb = DARK_TEXT

    # ================= SLIDE 9: TARGET USERS & IMPACT =================
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide9, "8. Target Users & Social Impact", "SLIDE 09")

    # Left Column: User Personas
    left_box = slide9.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(5.6), Inches(5.6))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Target User Personas"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    p.space_after = Pt(10)

    points = [
        "• Female Solo Commuters: Women navigating home late at night after work or university.",
        "• Night Shift Professionals: BPO, IT, and hospital staff commuting during night hours (8 PM - 5 AM).",
        "• Students & Tourists: Individuals unfamiliar with local crime hotspots in new cities."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(13.5)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)

    # Right Column: 3 Metric Cards
    metrics = [
        ("📉 70% Reduction in Unlit Exposure", "Reroutes pedestrians along brightly lit commercial arterial streets."),
        ("⚡ Sub-Second Emergency SOS Call Placement", "Placed directly to family via Twilio REST API without mobile CORS latency."),
        ("🗺️ 100% City-Wide Emergency Coverage", "Pre-loaded with 90+ verified emergency hubs across all zones of Hyderabad.")
    ]

    for j, (title, desc) in enumerate(metrics):
        add_flowchart_box(slide9, Inches(6.8), Inches(1.3 + j * 1.8), Inches(5.7), Inches(1.5), f"{title}\n\n{desc}", LIGHT_GREEN, EMERALD_GREEN, RGBColor(6, 95, 70), 12)

    # ================= SLIDE 10: FUTURE SCOPE & ROADMAP =================
    slide10 = prs.slides.add_slide(blank_layout)
    add_slide_header(slide10, "9. Future Scope & Roadmap", "SLIDE 10")

    # 3 Phase Horizontal Cards
    phases = [
        ("Phase 1: Pan-India Expansion", "• Scaling OpenStreetMap datasets to Mumbai, Delhi, Bengaluru, Chennai.\n• Crowd hazard report moderation."),
        ("Phase 2: IoT Smart Lighting", "• Direct API integration with municipal smart streetlight sensors.\n• Dynamic light intensity updates."),
        ("Phase 3: Wearable Trigger", "• Bluetooth smartwatch integration for silent 1-tap SOS trigger.\n• Automatic fall detection SOS.")
    ]

    for j, (p_title, p_desc) in enumerate(phases):
        left_pos = Inches(0.8 + j * 3.95)
        add_flowchart_box(slide10, left_pos, Inches(1.4), Inches(3.8), Inches(3.2), f"{p_title}\n\n{p_desc}", LIGHT_BLUE, SKY_BLUE, DARK_BLUE, 12)

    # Bottom Conclusion Card
    concl_card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.733), Inches(2.2))
    concl_card.fill.solid()
    concl_card.fill.fore_color.rgb = LIGHT_GREEN
    concl_card.line.color.rgb = EMERALD_GREEN
    concl_card.line.width = Pt(1.5)

    tf = concl_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🟢 SafeRoute: Transforming Passive Mapping into Active Safety Shields"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(6, 95, 70)
    p.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = "• Live Deployed Web Application: saferoute-tawny.vercel.app\n• GitHub Codebase: github.com/pranav-3010/saferoute\n\nThank you judges! We welcome your questions."
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = DARK_TEXT

    ppt_filename = "/Users/pranav/Desktop/saferoute/SafeRoute_SIH_10_Slide_Presentation.pptx"
    prs.save(ppt_filename)
    print("Enhanced PowerPoint Presentation Successfully Saved at:", ppt_filename)

if __name__ == '__main__':
    create_enhanced_presentation()
